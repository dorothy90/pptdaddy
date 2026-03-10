"""
Native PPTX export - Creates editable PowerPoint files with real text boxes, shapes, and tables.
No screenshots involved - all elements are native python-pptx objects.
"""

from pathlib import Path
from typing import Dict, List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


def _parse_color(hex_str: str) -> RGBColor:
    """Convert hex color string to RGBColor"""
    hex_str = hex_str.lstrip('#')
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


def _get_alignment(align_str: str) -> int:
    """Convert alignment string to PP_ALIGN enum"""
    mapping = {
        'left': PP_ALIGN.LEFT,
        'center': PP_ALIGN.CENTER,
        'right': PP_ALIGN.RIGHT,
        'justify': PP_ALIGN.JUSTIFY,
    }
    return mapping.get(align_str, PP_ALIGN.LEFT)


def _apply_font_style(run, style: Dict):
    """Apply font styling to a text run"""
    font = run.font
    if 'font_name' in style:
        font.name = style['font_name']
    if 'font_size' in style:
        font.size = Pt(style['font_size'])
    if 'font_bold' in style:
        font.bold = style['font_bold']
    if 'font_italic' in style:
        font.italic = style['font_italic']
    if 'font_color' in style:
        font.color.rgb = _parse_color(style['font_color'])


def _add_text_box(slide, element: Dict, prs):
    """Add a text box element to a slide"""
    pos = element['position']
    left = Inches(pos['left'])
    top = Inches(pos['top'])
    width = Inches(pos['width'])
    height = Inches(pos['height'])
    style = element.get('style', {})

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    # Vertical alignment
    v_align = style.get('vertical_alignment', 'top')
    if v_align == 'middle':
        tf.paragraphs[0].alignment = _get_alignment(style.get('alignment', 'left'))

    # Handle content - can be string or list of paragraphs
    content = element.get('content', '')

    if isinstance(content, list):
        # Multiple paragraphs / bullet points
        for i, item in enumerate(content):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            if isinstance(item, dict):
                p.text = item.get('text', '')
                item_style = {**style, **item.get('style', {})}
                _apply_font_style(p.runs[0] if p.runs else p.add_run(), item_style)
                if 'level' in item:
                    p.level = item['level']
            else:
                p.text = str(item)
                _apply_font_style(p.runs[0] if p.runs else p.add_run(), style)

            if 'alignment' in style:
                p.alignment = _get_alignment(style['alignment'])
            if 'line_spacing' in style:
                p.space_after = Pt(style['line_spacing'])
    else:
        # Single text
        tf.paragraphs[0].text = str(content)
        run = tf.paragraphs[0].runs[0] if tf.paragraphs[0].runs else tf.paragraphs[0].add_run()
        _apply_font_style(run, style)
        if 'alignment' in style:
            tf.paragraphs[0].alignment = _get_alignment(style['alignment'])

    # Background fill for text box
    if 'background_color' in style:
        fill = txBox.fill
        fill.solid()
        fill.fore_color.rgb = _parse_color(style['background_color'])


def _add_shape(slide, element: Dict, prs):
    """Add a shape element to a slide"""
    pos = element['position']
    left = Inches(pos['left'])
    top = Inches(pos['top'])
    width = Inches(pos['width'])
    height = Inches(pos['height'])
    style = element.get('style', {})

    shape_type_str = element.get('shape_type', 'rectangle')
    shape_map = {
        'rectangle': MSO_SHAPE.RECTANGLE,
        'rounded_rectangle': MSO_SHAPE.ROUNDED_RECTANGLE,
        'oval': MSO_SHAPE.OVAL,
        'circle': MSO_SHAPE.OVAL,
        'triangle': MSO_SHAPE.ISOSCELES_TRIANGLE,
        'diamond': MSO_SHAPE.DIAMOND,
        'chevron': MSO_SHAPE.CHEVRON,
        'arrow_right': MSO_SHAPE.RIGHT_ARROW,
        'arrow_left': MSO_SHAPE.LEFT_ARROW,
    }
    shape_type = shape_map.get(shape_type_str, MSO_SHAPE.RECTANGLE)

    shape = slide.shapes.add_shape(shape_type, left, top, width, height)

    # Fill color
    if 'fill_color' in style:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _parse_color(style['fill_color'])
    elif 'no_fill' in style and style['no_fill']:
        shape.fill.background()

    # Line/border
    if 'line_color' in style and style['line_color']:
        shape.line.color.rgb = _parse_color(style['line_color'])
        if 'line_width' in style:
            shape.line.width = Pt(style['line_width'])
    else:
        shape.line.fill.background()

    # Text inside shape
    if 'text' in element:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = element['text']
        text_style = style.get('text_style', style)
        run = tf.paragraphs[0].runs[0] if tf.paragraphs[0].runs else tf.paragraphs[0].add_run()
        _apply_font_style(run, text_style)
        if 'alignment' in text_style:
            tf.paragraphs[0].alignment = _get_alignment(text_style['alignment'])


def _add_table(slide, element: Dict, prs):
    """Add a table element to a slide"""
    pos = element['position']
    left = Inches(pos['left'])
    top = Inches(pos['top'])
    width = Inches(pos['width'])
    height = Inches(pos['height'])
    style = element.get('style', {})

    data = element.get('data', [])
    if not data:
        return

    rows = len(data)
    cols = len(data[0]) if data else 1

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    for row_idx, row_data in enumerate(data):
        for col_idx, cell_value in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_value)

            # Style header row
            if row_idx == 0 and 'header_color' in style:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _parse_color(style['header_color'])
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.bold = True
                        if 'header_font_color' in style:
                            run.font.color.rgb = _parse_color(style['header_font_color'])

            # Font size
            font_size = style.get('font_size', 14)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(font_size)


def _add_image(slide, element: Dict, prs):
    """Add an image element to a slide"""
    pos = element['position']
    left = Inches(pos['left'])
    top = Inches(pos['top'])
    width = Inches(pos['width'])
    height = Inches(pos['height'])

    image_path = element.get('image_path', '')
    if image_path and Path(image_path).exists():
        slide.shapes.add_picture(image_path, left, top, width, height)


def _set_background(slide, bg_def: Dict):
    """Set slide background color"""
    if not bg_def:
        return

    if 'color' in bg_def:
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = _parse_color(bg_def['color'])


def _add_slide(prs, slide_def: Dict):
    """Process one slide definition and add to presentation"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Set background
    if 'background' in slide_def:
        _set_background(slide, slide_def['background'])

    # Process elements
    element_handlers = {
        'text_box': _add_text_box,
        'shape': _add_shape,
        'table': _add_table,
        'image': _add_image,
    }

    for element in slide_def.get('elements', []):
        elem_type = element.get('type', '')
        handler = element_handlers.get(elem_type)
        if handler:
            try:
                handler(slide, element, prs)
            except Exception as e:
                print(f"   Warning: Failed to add {elem_type}: {e}")


def create_native_pptx(
    slides_data: List[Dict],
    output_file: str = "exports/presentation.pptx",
    presentation_title: str = "Presentation",
    progress_callback=None
) -> str:
    """
    Create a native editable PowerPoint file from structured slide data.

    Args:
        slides_data: List of slide definition dicts
        output_file: Output PPTX file path
        presentation_title: Title of the presentation
        progress_callback: Optional callback function(event_type, data)

    Returns:
        Path to the created PPTX file
    """
    print(f"\nCreating native PPTX presentation...")

    output_path = Path(output_file)
    output_path.parent.mkdir(exist_ok=True, parents=True)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    for i, slide_def in enumerate(slides_data, 1):
        try:
            _add_slide(prs, slide_def)
            print(f"   Added slide {i}")

            if progress_callback:
                progress_callback('pptx_slide_added', {
                    'slide_number': i,
                    'total_slides': len(slides_data)
                })
        except Exception as e:
            print(f"   Error adding slide {i}: {e}")

    prs.save(str(output_path))

    print(f"\nNative PPTX created: {output_path}")
    print(f"   Total slides: {len(prs.slides)}")
    print(f"   File size: {output_path.stat().st_size / 1024:.1f} KB\n")

    return str(output_path)
